#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Minimal RAG with Ollama (local LLM) + Chroma
Run:
  1) python -m venv .venv && source .venv/bin/activate  (Windows: .venv\\Scripts\\activate)
  2) pip install -r requirements.txt
  3) Make sure Ollama is installed & running: https://ollama.com
     - Pull a model once:  ollama pull llama3
  4) python main.py
"""
import os, glob, sys, shutil
from typing import List
from langchain.text_splitter import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter
from langchain.schema import Document
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.llms import Ollama
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA

DOC_DIR = os.environ.get("DOC_DIR", "docs")
MODEL_NAME = os.environ.get("OLLAMA_MODEL", "llama3")  # e.g., llama3, qwen2, mistral
CHUNK_SIZE = int(os.environ.get("CHUNK_SIZE", "500"))
CHUNK_OVERLAP = int(os.environ.get("CHUNK_OVERLAP", "75"))
TOP_K = int(os.environ.get("TOP_K", "5"))
TEMPERATURE = float(os.environ.get("TEMPERATURE", "0.7"))
SIMILARITY_THRESHOLD = float(os.environ.get("SIMILARITY_THRESHOLD", "10.0"))  # 관대한 임계값 (10.0 이하만 관련 있다고 판단)  # Chroma cosine distance: 낮을수록 유사
EMBED_MODEL = os.environ.get("EMBED_MODEL", "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2")

def load_raw_texts(doc_dir: str) -> List[str]:
    texts = []
    for path in sorted(glob.glob(os.path.join(doc_dir, "**", "*.*"), recursive=True)):
        filename = os.path.basename(path)
        print(f"[INFO] 처리 중: {filename}")
        
        if path.lower().endswith((".txt", ".md")):
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                    texts.append(f"파일명: {filename}\n내용:\n{content}")
            except Exception as e:
                print(f"[WARN] 텍스트 파일 읽기 실패: {path} ({e})")
                
        elif path.lower().endswith(".pdf"):
            try:
                from pypdf import PdfReader
                reader = PdfReader(path)
                content = "\n".join([p.extract_text() or "" for p in reader.pages])
                texts.append(f"파일명: {filename}\n내용:\n{content}")
            except Exception as e:
                print(f"[WARN] PDF 읽기 실패: {path} ({e})")
                
        elif path.lower().endswith(".docx"):
            try:
                from docx import Document
                doc = Document(path)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                texts.append(f"파일명: {filename}\n내용:\n{content}")
            except Exception as e:
                print(f"[WARN] Word 문서 읽기 실패: {path} ({e})")
                
        elif path.lower().endswith((".xlsx", ".xls")):
            try:
                import pandas as pd
                # 모든 시트 읽기
                excel_file = pd.ExcelFile(path)
                content_parts = []
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(path, sheet_name=sheet_name)
                    content_parts.append(f"시트: {sheet_name}\n{df.to_string()}")
                content = "\n\n".join(content_parts)
                texts.append(f"파일명: {filename}\n내용:\n{content}")
            except Exception as e:
                print(f"[WARN] Excel 파일 읽기 실패: {path} ({e})")
                
        elif path.lower().endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(path)
                content_parts = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    content_parts.append(f"슬라이드 {i}: {' '.join(slide_text)}")
                content = "\n\n".join(content_parts)
                texts.append(f"파일명: {filename}\n내용:\n{content}")
            except Exception as e:
                print(f"[WARN] PowerPoint 파일 읽기 실패: {path} ({e})")
                
    print(f"[INFO] 총 {len(texts)}개 파일 로드 완료")
    return texts

def smart_text_splitting(text: str, filename: str) -> List[Document]:
    """파일 타입에 따라 적절한 방법으로 텍스트 분할"""
    docs = []
    
    # 파일명에서 확장자 추출
    is_markdown = filename.lower().endswith('.md')
    
    if is_markdown:
        # 마크다운 파일: 헤더 기반 분할
        print(f"[INFO] 마크다운 헤더 기반 분할 적용: {filename}")
        
        # 마크다운 헤더 분할기 설정
        headers_to_split_on = [
            ("#", "Header 1"),
            ("##", "Header 2"), 
            ("###", "Header 3"),
            ("####", "Header 4"),
        ]
        
        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
        md_header_splits = markdown_splitter.split_text(text)
        
        # 각 헤더 섹션을 추가로 크기 기반으로 분할 (너무 큰 경우)
        for doc in md_header_splits:
            if len(doc.page_content) > CHUNK_SIZE * 2:  # 큰 섹션은 추가 분할
                sub_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=CHUNK_SIZE,
                    chunk_overlap=CHUNK_OVERLAP,
                    separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
                )
                sub_docs = sub_splitter.create_documents([doc.page_content])
                for sub_doc in sub_docs:
                    # 원래 헤더 메타데이터 유지
                    sub_doc.metadata.update(doc.metadata)
                    sub_doc.metadata['filename'] = filename
                    docs.append(sub_doc)
            else:
                doc.metadata['filename'] = filename
                docs.append(doc)
    
    else:
        # 일반 텍스트: 문단 우선 분할
        print(f"[INFO] 문단 기반 분할 적용: {filename}")
        
        # 문단 기반 분할기 (문단을 우선적으로 유지)
        paragraph_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", "! ", "? ", ", ", " ", ""]  # 문단(\n\n) 최우선
        )
        
        para_docs = paragraph_splitter.create_documents([text])
        for doc in para_docs:
            doc.metadata['filename'] = filename
            docs.append(doc)
    
    return docs

def build_retriever(texts: List[str]):
    
    # 이전 Chroma DB 삭제 (새로운 설정으로 재구성)
    chroma_dir = "./chroma_db"
    if os.path.exists(chroma_dir):
        shutil.rmtree(chroma_dir)
    
    docs = []
    for text in texts:
        # 파일명 추출 (파일명: {filename}\n내용:\n{content} 형태)
        lines = text.split('\n', 2)
        if len(lines) >= 3 and lines[0].startswith('파일명: '):
            filename = lines[0].replace('파일명: ', '')
            content = lines[2] if len(lines) > 2 else ''
        else:
            filename = "unknown.txt"
            content = text
        
        # 스마트 분할 적용
        file_docs = smart_text_splitting(content, filename)
        docs.extend(file_docs)

    print(f"[INFO] 총 {len(docs)}개 청크로 분할됨 (스마트 분할 적용)")
    
    # 분할 결과 상세 출력
    markdown_count = sum(1 for doc in docs if doc.metadata.get('filename', '').endswith('.md'))
    text_count = len(docs) - markdown_count
    print(f"[INFO] - 마크다운 청크: {markdown_count}개")
    print(f"[INFO] - 일반 텍스트 청크: {text_count}개")

    embed = HuggingFaceEmbeddings(model_name=EMBED_MODEL)
    vectordb = Chroma.from_documents(docs, embed, persist_directory=chroma_dir)
    
    # MMR(Maximal Marginal Relevance) 사용으로 다양성과 관련성 균형
    retriever = vectordb.as_retriever(
        search_type="mmr",
        search_kwargs={"k": TOP_K, "fetch_k": TOP_K * 2, "lambda_mult": 0.8}
    )
    return retriever

PROMPT_TMPL = """당신은 한국어로만 답변하는 AI 어시스턴트입니다.
다음 컨텍스트를 참고해서 질문에 반드시 한국어로만 답변해 주세요. 영어는 절대 사용하지 마세요.

컨텍스트:
{context}

질문: {question}

한국어 답변:"""

def build_chain(retriever):
    prompt = PromptTemplate(template=PROMPT_TMPL, input_variables=["context","question"])
    llm = Ollama(model=MODEL_NAME, temperature=TEMPERATURE)
    qa = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type="stuff",
        chain_type_kwargs={"prompt": prompt},
        return_source_documents=True,
    )
    return qa, llm

def is_relevant_to_docs(question: str, retriever, threshold: float = None):
    """질문이 문서와 관련이 있는지 유사도 점수로 자동 판단"""
    if threshold is None:
        threshold = SIMILARITY_THRESHOLD
        
    try:
        # 벡터 데이터베이스에서 유사도 점수와 함께 검색
        vectordb = retriever.vectorstore
        docs_and_scores = vectordb.similarity_search_with_score(question, k=TOP_K)
        
        if not docs_and_scores:
            return False, []
        
        print(f"[DEBUG] 질문: '{question}' | 검색 결과 유사도 점수:")
        for i, (doc, score) in enumerate(docs_and_scores):
            print(f"  {i+1}. 점수: {score:.3f} | 내용: {doc.page_content[:50]}...")
        print(f"[DEBUG] 임계값: {threshold}, 절대기준: 0.6")
        
        # Chroma cosine distance: 0(완전 유사) ~ 2(완전 다름)
        # 1. 절대적 안전 장치: 최고 점수가 0.6 이상이면 무조건 관련없음
        # 2. 가장 유사한 문서의 점수가 임계값보다 낮아야 함
        # 3. 상위 문서들의 평균 점수도 고려
        best_score = docs_and_scores[0][1]
        avg_score = sum(score for _, score in docs_and_scores[:3]) / min(3, len(docs_and_scores))
        
        # 절대적 안전 장치: 최고 점수가 15.0 이상이면 무조건 관련없음 (매우 관대하게)
        if best_score >= 15.0:
            print(f"[INFO] 절대적 기준으로 관련성 없음 - 최고 점수: {best_score:.3f} >= 15.0")
            return False, []
        
        if best_score <= threshold or avg_score <= threshold + 2.0:  # 훨씬 관대한 기준 적용
            # 임계값을 넘는 문서는 제외하되, 더 관대하게 판단
            relevant_docs = [doc for doc, score in docs_and_scores if score <= threshold + 5.0]  # 관대한 마진 적용
            if len(relevant_docs) >= 1:  # 최소 1개 이상의 관련 문서가 있어야 함
                print(f"[INFO] {len(relevant_docs)}개의 관련 문서 발견 (최고 점수: {best_score:.3f}, 평균: {avg_score:.3f}, 임계값: {threshold})")
                return True, relevant_docs
            
        print(f"[INFO] 관련성 낮음 - 최고: {best_score:.3f}, 평균: {avg_score:.3f} > 임계값: {threshold}")
        return False, []
        
        print(f"[INFO] 관련 문서 없음 - 최고 유사도: {docs_and_scores[0][1]:.3f} > 임계값: {threshold}")
        return False, []
        
    except Exception as e:
        print(f"[WARN] 관련성 판단 중 오류: {e}")
        return False, []

def main():
    if not os.path.isdir(DOC_DIR):
        print(f"[ERR] 문서 폴더가 없습니다: {DOC_DIR}")
        sys.exit(1)

    print(f"[INFO] 문서 인덱싱 중... ({DOC_DIR})")
    texts = load_raw_texts(DOC_DIR)
    if not texts:
        print("[ERR] 인덱싱할 문서가 없습니다. docs 폴더에 txt/md/pdf를 넣으세요.")
        sys.exit(1)

    retriever = build_retriever(texts)
    qa, llm = build_chain(retriever)
    print("[OK] RAG 준비 완료. 아래에 질문을 입력하세요. (종료: 빈 줄 + Enter)")

    while True:
        try:
            q = input("\n질문> ").strip()
            if not q:
                break
            
            # 질문이 문서와 관련이 있는지 확인
            is_relevant, relevant_docs = is_relevant_to_docs(q, retriever)
            
            if is_relevant:
                # RAG 사용 - 필터링된 관련 문서로만 답변 생성
                if relevant_docs:
                    # 관련 문서들을 컨텍스트로 조합
                    context = "\n\n".join([doc.page_content for doc in relevant_docs])
                    prompt = f"다음 컨텍스트를 참고해서 질문에 한국어로 답변해 주세요.\n\n컨텍스트:\n{context}\n\n질문: {q}\n\n답변:"
                    
                    # LLM으로 직접 답변 생성
                    response = llm.invoke(prompt)
                    print("\n--- 답변 (RAG) ---")
                    print(response)
                    
                    print("\n[근거 출처]")
                    for i, doc in enumerate(relevant_docs, 1):
                        snippet = doc.page_content.strip().replace("\n"," ")
                        if len(snippet) > 120:
                            snippet = snippet[:120] + "…"
                        print(f"{i}. {snippet}")
                else:
                    # 키워드는 있지만 관련 문서가 없는 경우
                    response = llm.invoke(f"다음 질문에 한국어로 답변해주세요: {q}")
                    print("\n--- 답변 (LLM) ---")
                    print(response)
            else:
                # 순수 LLM 사용
                response = llm.invoke(f"다음 질문에 한국어로 답변해주세요: {q}")
                print("\n--- 답변 (LLM) ---")
                print(response)
                
        except (KeyboardInterrupt, EOFError):
            break

if __name__ == "__main__":
    main()
